from pptx import Presentation 
from datetime import datetime
import os
from PIL import Image
import glob

from typing import Union
import json
import pandas as pd 
import requests
import urllib
import os 
from fake_useragent import UserAgent
from requests.exceptions import HTTPError

import gtts
from playsound import playsound

prompt = "Welcome to Presentation Mode! Here, we will get you inspired to start your presentation by generating a slide deck of free-to-use images from the Internet based on the options you give me! This mode works best if the presentation you want to make is about a visual topic. But feel free to make a presentation about any topic!"
print(prompt)
tts = gtts.gTTS(prompt)
tts.save("prompt.mp3")
playsound("prompt.mp3")


X = Presentation()
Layout = X.slide_layouts[0] 
first_slide = X.slides.add_slide(Layout) # Adding first slide
prompt = "Enter the title of your presentation: "
print(prompt)
tts = gtts.gTTS(prompt)
tts.save("prompt.mp3")
playsound("prompt.mp3")
input_str = str(input())
first_slide.shapes.title.text = input_str
presentation_title = input_str
prompt = "Enter the name of presenters:"
print(prompt)
tts = gtts.gTTS(prompt)
tts.save("prompt.mp3")
playsound("prompt.mp3")
author_names = str(input())
first_slide.placeholders[1].text = "By " + author_names
now = datetime.now()
new_now = now.strftime("%H%M%S")


def call_request(url) -> Union[HTTPError, dict]:
    user_agent = UserAgent()
    headers = headers={'User-Agent': str(user_agent)}
    response = requests.get(url, headers=headers)
    try:
        response.raise_for_status()
    except requests.exceptions.HTTPError as e:
        return e

    return response.json()

prompt = "Enter 2-3 keywords (separated by spaces): "
print(prompt)
tts = gtts.gTTS(prompt)
tts.save("prompt.mp3")
playsound("prompt.mp3")
keyword = str(input())
prompt = "Creating presentation now. Please feel free to start imagining what your presentation will look like while I create a starter presentation!"
print(prompt)
tts = gtts.gTTS(prompt)
tts.save("prompt.mp3")
playsound("prompt.mp3")
genre= keyword
per_page=20
page=1
image_folder_path= os.getcwd()+"/images"
if not os.path.isdir(image_folder_path):
    os.mkdir(image_folder_path)
parameter={"query":genre,"per_page":per_page,"page":page}
query= urllib.parse.urlencode(parameter)
url=f"https://unsplash.com/napi/search/photos?{query}"
response=call_request(url)
image_list=[]
if len(response['results'])>0:
    for i in range(len(response['results'])):
        filename = response['results'][i]['urls']['raw'].split('/')[-1].split('?')[0]+".jpg"
        folder_path=os.path.join(image_folder_path,genre)
        if not os.path.isdir(folder_path):
            os.mkdir(folder_path)
        filepath=os.path.join(folder_path,filename)
        r = requests.get(response['results'][i]['urls']['raw'], allow_redirects=True)
        open(filepath.replace("\\", "/"), 'wb').write(r.content)
        temp={ "Genre":genre, "link":response['results'][i]['urls']['raw']}
        image_list.append(temp)
    





# second_Layout = X.slide_layouts[8]
# second_slide = X.slides.add_slide(second_Layout)
# placeholder = second_slide.placeholders[1]
# picture = placeholder.insert_picture('UserAmplitudePlot.png')
full_directory = image_folder_path + '/' + keyword + "/*.jpg"
print("full directory", full_directory)
# full_directory = './images/parisfood/*.jpg'
# full_directory = 'C:/Users/ishan/Desktop/2022Spring/6.835/Project/GazeTracking/images/parisfood'
for filename in glob.glob(full_directory):
	try:
		# print(filename)
		x_Layout = X.slide_layouts[8]
		x_slide = X.slides.add_slide(x_Layout)
		placeholder = x_slide.placeholders[1]
		# im = Image.open(filename)
		picture = placeholder.insert_picture(filename)
	except:
		pass

presentation_name = presentation_title + str(new_now) + ".pptx"
X.save(presentation_name)
print("Starter presentation file successfully created! To see your starter presentation, please navigate to the Carvis folder and see the file with the following name:", presentation_name)
# os.startfile(presentation_name)

print('To try a different mode, please enter the following command: python3 cli.py')





 



